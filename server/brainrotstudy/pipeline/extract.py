"""Stage 1: extract normalized content from a topic, PDF, PPTX, or text file."""

from __future__ import annotations

import re
from pathlib import Path

from ..schemas import ExtractedContent, ExtractedSection


MAX_SECTIONS = 20
MAX_BULLETS_PER_SECTION = 8
MAX_BODY_CHARS_PER_SECTION = 1200


def extract_from_topic(topic: str, outline: str | None = None) -> ExtractedContent:
    title = topic.strip().splitlines()[0][:120] if topic else "Study Session"
    sections: list[ExtractedSection] = []
    if outline:
        sections = _sections_from_outline(outline)
    return ExtractedContent(title=title, source="topic", sections=sections)


def extract_from_path(path: Path) -> ExtractedContent:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf(path)
    if suffix == ".pptx":
        return _extract_pptx(path)
    if suffix in {".txt", ".md"}:
        return _extract_text(path)
    raise ValueError(f"Unsupported file type: {suffix}")


# ----- PDF -----------------------------------------------------------------


def _extract_pdf(path: Path) -> ExtractedContent:
    try:
        import pymupdf  # type: ignore
    except ImportError as e:
        raise RuntimeError(
            "PDF extraction requires PyMuPDF. Install with: uv pip install pymupdf"
        ) from e

    title = path.stem.replace("_", " ").replace("-", " ").title()
    sections: list[ExtractedSection] = []
    with pymupdf.open(path) as doc:
        if doc.metadata and doc.metadata.get("title"):
            title = doc.metadata["title"][:120]
        for page_idx, page in enumerate(doc):
            if len(sections) >= MAX_SECTIONS:
                break
            text = page.get_text("text").strip()
            if not text:
                continue
            section = _section_from_block(text, default_heading=f"Page {page_idx + 1}")
            if section.bullets or section.body:
                sections.append(section)
    return ExtractedContent(title=title, source="pdf", sections=sections)


# ----- PPTX ----------------------------------------------------------------


def _extract_pptx(path: Path) -> ExtractedContent:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as e:
        raise RuntimeError(
            "PPTX extraction requires python-pptx. Install with: uv pip install python-pptx"
        ) from e

    prs = Presentation(str(path))
    title = path.stem.replace("_", " ").replace("-", " ").title()
    sections: list[ExtractedSection] = []
    for idx, slide in enumerate(prs.slides):
        if len(sections) >= MAX_SECTIONS:
            break
        heading = ""
        bullets: list[str] = []
        body_parts: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = "".join(run.text for run in paragraph.runs).strip()
                if not text:
                    continue
                if not heading and len(text) < 100 and paragraph.level == 0:
                    heading = text
                elif paragraph.level > 0 or len(text) < 200:
                    bullets.append(text)
                else:
                    body_parts.append(text)
        section = ExtractedSection(
            heading=heading or f"Slide {idx + 1}",
            bullets=bullets[:MAX_BULLETS_PER_SECTION],
            body=" ".join(body_parts)[:MAX_BODY_CHARS_PER_SECTION],
        )
        if section.bullets or section.body:
            sections.append(section)
    if idx == 0 and title == path.stem.replace("_", " ").replace("-", " ").title():
        if sections and sections[0].heading:
            title = sections[0].heading
    return ExtractedContent(title=title, source="pptx", sections=sections)


# ----- Plain text / Markdown ----------------------------------------------


def _extract_text(path: Path) -> ExtractedContent:
    raw = path.read_text(encoding="utf-8", errors="replace")
    lines = raw.splitlines()
    title = path.stem.replace("_", " ").replace("-", " ").title()
    for line in lines:
        stripped = line.strip()
        if stripped.startswith("# "):
            title = stripped[2:].strip()[:120]
            break
        if stripped:
            title = stripped[:120]
            break

    sections = _sections_from_outline(raw)
    return ExtractedContent(title=title, source="text", sections=sections)


# ----- Helpers -------------------------------------------------------------


_HEADING_RE = re.compile(r"^(#{1,6})\s+(.+)$")
_BULLET_RE = re.compile(r"^\s*[-*•]\s+(.+)$")


def _sections_from_outline(text: str) -> list[ExtractedSection]:
    """Split an outline (markdown-ish) into heading + bullets + body sections."""
    sections: list[ExtractedSection] = []
    current = ExtractedSection()
    paragraph_buf: list[str] = []

    def flush_paragraph() -> None:
        if paragraph_buf:
            current.body = (current.body + " " + " ".join(paragraph_buf)).strip()
            paragraph_buf.clear()

    def start_new(heading: str) -> None:
        nonlocal current
        flush_paragraph()
        if current.heading or current.bullets or current.body:
            sections.append(current)
        current = ExtractedSection(heading=heading)

    for raw_line in text.splitlines():
        line = raw_line.rstrip()
        if not line.strip():
            flush_paragraph()
            continue
        heading_match = _HEADING_RE.match(line)
        if heading_match:
            start_new(heading_match.group(2).strip()[:120])
            continue
        bullet_match = _BULLET_RE.match(line)
        if bullet_match:
            flush_paragraph()
            if len(current.bullets) < MAX_BULLETS_PER_SECTION:
                current.bullets.append(bullet_match.group(1).strip())
            continue
        paragraph_buf.append(line.strip())

    flush_paragraph()
    if current.heading or current.bullets or current.body:
        sections.append(current)

    for s in sections:
        s.body = s.body[:MAX_BODY_CHARS_PER_SECTION]

    return sections[:MAX_SECTIONS]


def _section_from_block(text: str, *, default_heading: str) -> ExtractedSection:
    """Best-effort split of free text into a heading + body."""
    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
    if not lines:
        return ExtractedSection(heading=default_heading)

    heading = lines[0][:120] if len(lines[0]) < 120 else default_heading
    start_body_idx = 1 if heading != default_heading else 0

    bullets: list[str] = []
    body_parts: list[str] = []
    for ln in lines[start_body_idx:]:
        bullet = _BULLET_RE.match(ln)
        if bullet and len(bullets) < MAX_BULLETS_PER_SECTION:
            bullets.append(bullet.group(1).strip())
        else:
            body_parts.append(ln)

    body = " ".join(body_parts)[:MAX_BODY_CHARS_PER_SECTION]
    return ExtractedSection(heading=heading, bullets=bullets, body=body)
