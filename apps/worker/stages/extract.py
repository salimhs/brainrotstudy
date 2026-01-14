"""Stage 1: Extract content from PDF/PPTX/DOCX/XLSX/TXT/MD/Images."""

import json
import sys
from pathlib import Path
from typing import List, Optional

sys.path.insert(0, "/app")

from shared.models import SlidesExtracted, SlideData
from shared.utils import get_job_dir, get_job_logger, load_job_metadata, artifact_exists


# Supported file extensions by category
DOCUMENT_EXTENSIONS = {".pdf", ".pptx", ".docx"}
SPREADSHEET_EXTENSIONS = {".xlsx", ".xls", ".csv"}
TEXT_EXTENSIONS = {".txt", ".md", ".markdown"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"}

ALL_EXTENSIONS = DOCUMENT_EXTENSIONS | SPREADSHEET_EXTENSIONS | TEXT_EXTENSIONS | IMAGE_EXTENSIONS


def run_extract_stage(job_id: str) -> None:
    """
    Extract text and images from uploaded files.
    
    Supports: PDF, PPTX, DOCX, XLSX/CSV, TXT, Markdown, Images (with OCR).
    For topic-only jobs, this stage is skipped.
    """
    logger = get_job_logger(job_id)
    job_dir = get_job_dir(job_id)
    metadata = load_job_metadata(job_id)
    
    if not metadata:
        raise ValueError(f"Job {job_id} not found")
    
    # Ensure extracted directory exists
    extracted_dir = job_dir / "extracted"
    extracted_dir.mkdir(parents=True, exist_ok=True)
    
    # Check for idempotency
    slides_json_path = extracted_dir / "slides.json"
    if slides_json_path.exists():
        logger.info("CACHE HIT: slides.json already exists, skipping extraction")
        return
    
    # For topic-only jobs, create empty slides structure
    if metadata.input_type == "topic":
        logger.info("Topic-only job, creating empty slides structure")
        slides = SlidesExtracted(slides=[])
        with open(slides_json_path, "w", encoding="utf-8") as f:
            json.dump(slides.model_dump(), f, indent=2)
        return
    
    # Find input files with all supported extensions
    input_dir = job_dir / "input"
    input_files = []
    for ext in ALL_EXTENSIONS:
        input_files.extend(input_dir.glob(f"*{ext}"))
    
    if not input_files:
        logger.warning("No input file found, creating empty slides structure")
        slides = SlidesExtracted(slides=[])
        with open(slides_json_path, "w", encoding="utf-8") as f:
            json.dump(slides.model_dump(), f, indent=2)
        return
    
    # Process the first input file (could be extended to merge multiple files)
    input_file = input_files[0]
    ext = input_file.suffix.lower()
    logger.info(f"Processing file: {input_file.name} (type: {ext})")
    
    # Route to appropriate extractor
    if ext == ".pdf":
        slides = extract_pdf(job_id, input_file)
    elif ext == ".pptx":
        slides = extract_pptx(job_id, input_file)
    elif ext == ".docx":
        slides = extract_docx(job_id, input_file)
    elif ext in SPREADSHEET_EXTENSIONS:
        slides = extract_spreadsheet(job_id, input_file)
    elif ext in TEXT_EXTENSIONS:
        slides = extract_text(job_id, input_file)
    elif ext in IMAGE_EXTENSIONS:
        slides = extract_image(job_id, input_file)
    else:
        raise ValueError(f"Unsupported file type: {ext}")
    
    # Save extracted slides
    with open(slides_json_path, "w", encoding="utf-8") as f:
        json.dump(slides.model_dump(), f, indent=2)
    
    logger.info(f"Extracted {len(slides.slides)} slides/sections")


def extract_pdf(job_id: str, pdf_path: Path) -> SlidesExtracted:
    """Extract content from a PDF file using PyMuPDF."""
    logger = get_job_logger(job_id)
    job_dir = get_job_dir(job_id)
    extracted_dir = job_dir / "extracted"
    extracted_dir.mkdir(parents=True, exist_ok=True)
    
    slides = []
    
    try:
        import fitz  # PyMuPDF
        
        doc = fitz.open(pdf_path)
        
        for page_num, page in enumerate(doc, start=1):
            logger.info(f"Processing PDF page {page_num}")
            
            # Extract text
            text = page.get_text()
            
            # Render page as image
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x scale for quality
            image_path = extracted_dir / f"slide_{page_num:02d}.png"
            pix.save(str(image_path))
            
            # Parse title (first line) and bullets
            lines = [line.strip() for line in text.split('\n') if line.strip()]
            title = lines[0] if lines else f"Slide {page_num}"
            bullets = lines[1:] if len(lines) > 1 else []
            
            slides.append(SlideData(
                index=page_num,
                title=title,
                bullets=bullets[:10],  # Limit bullets
                raw_text=text,
                rendered_image=f"extracted/slide_{page_num:02d}.png",
                images=[],
            ))
        
        doc.close()
        
    except ImportError:
        logger.warning("PyMuPDF not installed, creating placeholder slide")
        slides.append(SlideData(
            index=1,
            title="PDF Content",
            bullets=["Content extracted from PDF"],
            raw_text="PDF content placeholder",
        ))
    except Exception as e:
        logger.error(f"PDF extraction error: {e}")
        slides.append(SlideData(
            index=1,
            title="PDF Content",
            bullets=[f"Error extracting: {str(e)}"],
            raw_text="",
        ))
    
    return SlidesExtracted(slides=slides)


def extract_pptx(job_id: str, pptx_path: Path) -> SlidesExtracted:
    """Extract content from a PPTX file using python-pptx."""
    logger = get_job_logger(job_id)
    job_dir = get_job_dir(job_id)
    extracted_dir = job_dir / "extracted"
    extracted_dir.mkdir(parents=True, exist_ok=True)
    
    slides = []
    
    try:
        from pptx import Presentation
        
        prs = Presentation(pptx_path)
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            logger.info(f"Processing PPTX slide {slide_num}")
            
            title = ""
            bullets = []
            raw_text_parts = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        raw_text_parts.append(text)
                        
                        # First text shape is usually title
                        if not title and hasattr(shape, "is_placeholder"):
                            title = text
                        else:
                            # Split into bullets
                            for line in text.split('\n'):
                                if line.strip():
                                    bullets.append(line.strip())
            
            if not title:
                title = f"Slide {slide_num}"
            
            slides.append(SlideData(
                index=slide_num,
                title=title,
                bullets=bullets[:10],
                raw_text="\n".join(raw_text_parts),
                images=[],
            ))
        
    except ImportError:
        logger.warning("python-pptx not installed, creating placeholder slide")
        slides.append(SlideData(
            index=1,
            title="PPTX Content",
            bullets=["Content extracted from PPTX"],
            raw_text="PPTX content placeholder",
        ))
    except Exception as e:
        logger.error(f"PPTX extraction error: {e}")
        slides.append(SlideData(
            index=1,
            title="PPTX Content",
            bullets=[f"Error extracting: {str(e)}"],
            raw_text="",
        ))
    
    return SlidesExtracted(slides=slides)


def extract_docx(job_id: str, docx_path: Path) -> SlidesExtracted:
    """Extract content from a DOCX file using python-docx."""
    logger = get_job_logger(job_id)
    slides = []
    
    try:
        from docx import Document
        
        doc = Document(docx_path)
        logger.info(f"Processing DOCX: {docx_path.name}")
        
        # Group paragraphs into sections based on headings
        current_section = {
            "title": "Document Content",
            "bullets": [],
            "raw_text_parts": []
        }
        sections = []
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            
            # Check if this is a heading
            if para.style.name.startswith('Heading'):
                # Save current section if it has content
                if current_section["bullets"] or current_section["raw_text_parts"]:
                    sections.append(current_section)
                # Start new section
                current_section = {
                    "title": text,
                    "bullets": [],
                    "raw_text_parts": []
                }
            else:
                current_section["raw_text_parts"].append(text)
                # Treat short paragraphs as bullets
                if len(text) < 200:
                    current_section["bullets"].append(text)
        
        # Don't forget the last section
        if current_section["bullets"] or current_section["raw_text_parts"]:
            sections.append(current_section)
        
        # Convert sections to slides
        for idx, section in enumerate(sections, start=1):
            slides.append(SlideData(
                index=idx,
                title=section["title"],
                bullets=section["bullets"][:10],
                raw_text="\n".join(section["raw_text_parts"]),
                images=[],
            ))
        
        # If no sections found, create a single slide with all content
        if not slides:
            all_text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
            lines = [line.strip() for line in all_text.split('\n') if line.strip()]
            slides.append(SlideData(
                index=1,
                title=lines[0] if lines else "Document Content",
                bullets=lines[1:11] if len(lines) > 1 else [],
                raw_text=all_text,
                images=[],
            ))
        
        logger.info(f"Extracted {len(slides)} sections from DOCX")
        
    except ImportError:
        logger.warning("python-docx not installed, creating placeholder")
        slides.append(SlideData(
            index=1,
            title="DOCX Content",
            bullets=["python-docx library not available"],
            raw_text="",
        ))
    except Exception as e:
        logger.error(f"DOCX extraction error: {e}")
        slides.append(SlideData(
            index=1,
            title="DOCX Content",
            bullets=[f"Error extracting: {str(e)}"],
            raw_text="",
        ))
    
    return SlidesExtracted(slides=slides)


def extract_spreadsheet(job_id: str, file_path: Path) -> SlidesExtracted:
    """Extract content from XLSX/XLS/CSV files."""
    logger = get_job_logger(job_id)
    slides = []
    ext = file_path.suffix.lower()
    
    try:
        if ext == ".csv":
            import csv
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                reader = csv.reader(f)
                rows = list(reader)
            
            if rows:
                headers = rows[0] if rows else []
                data_rows = rows[1:21]  # Limit to 20 rows
                
                # Create summary bullets
                bullets = [f"Columns: {', '.join(headers[:5])}{'...' if len(headers) > 5 else ''}"]
                bullets.append(f"Total rows: {len(rows) - 1}")
                for row in data_rows[:5]:
                    bullets.append(" | ".join(str(cell)[:30] for cell in row[:4]))
                
                slides.append(SlideData(
                    index=1,
                    title=f"Data from {file_path.name}",
                    bullets=bullets[:10],
                    raw_text="\n".join([",".join(row) for row in rows[:50]]),
                    images=[],
                ))
        else:
            # XLSX/XLS using openpyxl
            from openpyxl import load_workbook
            
            wb = load_workbook(file_path, read_only=True, data_only=True)
            
            for sheet_idx, sheet_name in enumerate(wb.sheetnames[:5], start=1):  # Limit to 5 sheets
                sheet = wb[sheet_name]
                logger.info(f"Processing sheet: {sheet_name}")
                
                rows = []
                for row_idx, row in enumerate(sheet.iter_rows(values_only=True)):
                    if row_idx >= 20:  # Limit rows
                        break
                    rows.append([str(cell) if cell is not None else "" for cell in row])
                
                if rows:
                    headers = rows[0] if rows else []
                    bullets = [f"Columns: {', '.join(headers[:5])}{'...' if len(headers) > 5 else ''}"]
                    bullets.append(f"Sheet: {sheet_name}")
                    for row in rows[1:6]:
                        bullets.append(" | ".join(str(cell)[:30] for cell in row[:4]))
                    
                    slides.append(SlideData(
                        index=sheet_idx,
                        title=f"Sheet: {sheet_name}",
                        bullets=bullets[:10],
                        raw_text="\n".join(["\t".join(row) for row in rows]),
                        images=[],
                    ))
            
            wb.close()
        
        if not slides:
            slides.append(SlideData(
                index=1,
                title="Spreadsheet Content",
                bullets=["No data found in spreadsheet"],
                raw_text="",
                images=[],
            ))
        
        logger.info(f"Extracted {len(slides)} sheets/sections from spreadsheet")
        
    except ImportError as e:
        logger.warning(f"Spreadsheet library not installed: {e}")
        slides.append(SlideData(
            index=1,
            title="Spreadsheet Content",
            bullets=["openpyxl library not available"],
            raw_text="",
        ))
    except Exception as e:
        logger.error(f"Spreadsheet extraction error: {e}")
        slides.append(SlideData(
            index=1,
            title="Spreadsheet Content",
            bullets=[f"Error extracting: {str(e)}"],
            raw_text="",
        ))
    
    return SlidesExtracted(slides=slides)


def extract_text(job_id: str, file_path: Path) -> SlidesExtracted:
    """Extract content from TXT or Markdown files."""
    logger = get_job_logger(job_id)
    slides = []
    ext = file_path.suffix.lower()
    
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        logger.info(f"Processing text file: {file_path.name} ({len(content)} chars)")
        
        # For Markdown, try to parse sections
        if ext in {".md", ".markdown"}:
            try:
                import markdown
                import re
                
                # Split by headers
                sections = re.split(r'\n#{1,3}\s+', content)
                header_matches = re.findall(r'\n(#{1,3})\s+(.+)', '\n' + content)
                
                if header_matches:
                    for idx, (level, title) in enumerate(header_matches, start=1):
                        section_content = sections[idx] if idx < len(sections) else ""
                        lines = [line.strip() for line in section_content.split('\n') if line.strip()]
                        
                        slides.append(SlideData(
                            index=idx,
                            title=title.strip(),
                            bullets=lines[:10],
                            raw_text=section_content,
                            images=[],
                        ))
                else:
                    # No headers, treat as single section
                    lines = [line.strip() for line in content.split('\n') if line.strip()]
                    slides.append(SlideData(
                        index=1,
                        title=lines[0][:80] if lines else "Text Content",
                        bullets=lines[1:11] if len(lines) > 1 else [],
                        raw_text=content,
                        images=[],
                    ))
                    
            except ImportError:
                # Fallback without markdown library
                lines = [line.strip() for line in content.split('\n') if line.strip()]
                slides.append(SlideData(
                    index=1,
                    title=lines[0][:80] if lines else "Markdown Content",
                    bullets=lines[1:11] if len(lines) > 1 else [],
                    raw_text=content,
                    images=[],
                ))
        else:
            # Plain text: split into paragraphs
            paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]
            
            if paragraphs:
                # Create slides from paragraphs (max 10 paragraphs per slide)
                for idx, para in enumerate(paragraphs[:10], start=1):
                    lines = [line.strip() for line in para.split('\n') if line.strip()]
                    slides.append(SlideData(
                        index=idx,
                        title=lines[0][:80] if lines else f"Section {idx}",
                        bullets=lines[1:] if len(lines) > 1 else [],
                        raw_text=para,
                        images=[],
                    ))
            else:
                lines = [line.strip() for line in content.split('\n') if line.strip()]
                slides.append(SlideData(
                    index=1,
                    title=lines[0][:80] if lines else "Text Content",
                    bullets=lines[1:11] if len(lines) > 1 else [],
                    raw_text=content,
                    images=[],
                ))
        
        logger.info(f"Extracted {len(slides)} sections from text file")
        
    except Exception as e:
        logger.error(f"Text extraction error: {e}")
        slides.append(SlideData(
            index=1,
            title="Text Content",
            bullets=[f"Error reading file: {str(e)}"],
            raw_text="",
        ))
    
    return SlidesExtracted(slides=slides)


def extract_image(job_id: str, image_path: Path) -> SlidesExtracted:
    """Extract text from images using OCR (pytesseract)."""
    logger = get_job_logger(job_id)
    job_dir = get_job_dir(job_id)
    extracted_dir = job_dir / "extracted"
    extracted_dir.mkdir(parents=True, exist_ok=True)
    
    slides = []
    
    try:
        from PIL import Image
        
        # Copy image to extracted dir
        img = Image.open(image_path)
        output_path = extracted_dir / f"slide_01{image_path.suffix}"
        img.save(output_path)
        
        logger.info(f"Processing image: {image_path.name}")
        
        # Try OCR extraction
        ocr_text = ""
        try:
            import pytesseract
            ocr_text = pytesseract.image_to_string(img)
            logger.info(f"OCR extracted {len(ocr_text)} characters")
        except ImportError:
            logger.warning("pytesseract not installed, skipping OCR")
        except Exception as e:
            logger.warning(f"OCR failed: {e}")
        
        # Parse OCR text into bullets
        lines = [line.strip() for line in ocr_text.split('\n') if line.strip()]
        
        slides.append(SlideData(
            index=1,
            title=lines[0][:80] if lines else f"Image: {image_path.name}",
            bullets=lines[1:11] if len(lines) > 1 else ["Image content"],
            raw_text=ocr_text or f"Image file: {image_path.name}",
            rendered_image=f"extracted/slide_01{image_path.suffix}",
            images=[],
        ))
        
    except ImportError:
        logger.warning("Pillow not installed")
        slides.append(SlideData(
            index=1,
            title=f"Image: {image_path.name}",
            bullets=["Image processing unavailable"],
            raw_text="",
        ))
    except Exception as e:
        logger.error(f"Image extraction error: {e}")
        slides.append(SlideData(
            index=1,
            title=f"Image: {image_path.name}",
            bullets=[f"Error processing: {str(e)}"],
            raw_text="",
        ))
    
    return SlidesExtracted(slides=slides)
